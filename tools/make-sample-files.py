# -*- coding: utf-8 -*-
"""테스트용 샘플 파일 묶음 생성기 — 지원 형식별로 파일 하나씩 만들어 한 폴더에 담는다.

왜: 손으로 여러 형식의 파일을 모으려면 오피스·한글·인코더가 다 있어야 하고, 그렇게 모은
    바이너리를 저장소에 커밋하면 용량만 쌓인다. 그래서 "폴더째 열어 눈으로 확인하는" 용도의
    파일 묶음은 코드로 만들어 두고, 필요할 때 다시 생성한다(단위 테스트 픽스처인
    tests/fixtures/*.js 와 같은 생각. 구형 .doc 은 그 픽스처를 node 로 그대로 재사용한다).

사용법:
    python tools/make-sample-files.py [출력폴더]
    (기본 출력 위치: 저장소 바로 옆 "테스트파일모음" 폴더 — 저장소 안을 더럽히지 않는다)

필요한 것(없으면 그 형식만 건너뛰고 나머지는 계속 만든다):
    pip install reportlab openpyxl python-pptx pillow opencv-python
    node        — 구형 .doc (tests/fixtures/build-doc.js)
    맑은 고딕   — 이미지·영상 안의 한글 글자

만들 수 없어 빠지는 형식: .hwp(구형 한글)·.xls(구형 엑셀)은 바이너리 규격이라 한글/엑셀에서
직접 저장해야 하고, .m4a/.aac/.ogg/.flac 과 .mov/.mkv/.wmv/.flv 는 인코더가 필요하다.
"""
import base64, io, json, math, os, shutil, sqlite3, struct, subprocess, sys, tarfile, tempfile, time, wave, zipfile

# 콘솔 코드페이지(cp949 등)에 없는 글자 때문에 진행 상황 출력이 죽지 않도록.
try:
    sys.stdout.reconfigure(errors="replace")
except Exception:
    pass

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT = os.path.abspath(sys.argv[1]) if len(sys.argv) > 1 \
    else os.path.join(os.path.dirname(REPO), "테스트파일모음")
KFONT = r"C:\Windows\Fonts\malgun.ttf"
KFONTB = r"C:\Windows\Fonts\malgunbd.ttf"

made, skipped = [], []


def p(*parts):
    path = os.path.join(OUT, *parts)
    os.makedirs(os.path.dirname(path), exist_ok=True)
    return path


def write_text(name, text, bom=False):
    with open(p(name), "w", encoding="utf-8-sig" if bom else "utf-8", newline="\n") as f:
        f.write(text)


def done(name):
    made.append(name)
    print("  +", name)


def step(label, fn):
    """한 형식이 실패해도(라이브러리 없음 등) 나머지는 계속 만든다."""
    try:
        fn()
    except Exception as e:
        skipped.append("%s: %s" % (label, e))
        print("  - 건너뜀:", label, ":", e)


# ---------------------------------------------------------------- PDF
def make_pdf():
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    pdfmetrics.registerFont(UnicodeCIDFont("HYSMyeongJo-Medium"))   # 글꼴 내장 없이 쓰는 한글 CID 글꼴
    F = "HYSMyeongJo-Medium"
    W, H = A4
    c = canvas.Canvas(p("문서-샘플.pdf"), pagesize=A4)

    def page_no(n):
        c.setFont(F, 9); c.setFillColorRGB(.45, .45, .45)
        c.drawCentredString(W / 2, 30, "- %d -" % n); c.setFillColorRGB(0, 0, 0)

    c.setFont(F, 22); c.drawString(60, H - 90, "만능파일교실 PDF 테스트 문서")
    c.setFont(F, 11)
    lines = [
        "이 문서는 PDF 보기·편집 기능을 확인하려고 만든 샘플입니다.",
        "",
        "확인해 볼 것",
        "  1. 페이지 이동 / 확대·축소 / 회전",
        "  2. 찾기 — '검색테스트' 라는 낱말이 3쪽에 한 번 나옵니다.",
        "  3. 서명·텍스트·날짜·체크 넣고 저장하기",
        "  4. 펜·형광펜 필기, 책갈피 목차, 코드 핀",
        "  5. 페이지 정리(회전·삭제·순서 바꾸기)와 다른 PDF 와 합치기",
        "",
        "아래 빈칸에 서명과 날짜를 넣어 보세요.",
    ]
    y = H - 140
    for ln in lines:
        c.drawString(60, y, ln); y -= 20
    c.setFont(F, 12)
    c.drawString(60, 300, "서 명 :"); c.line(120, 296, 320, 296)
    c.drawString(60, 250, "날 짜 :"); c.line(120, 246, 320, 246)
    c.rect(60, 190, 14, 14); c.drawString(84, 192, "위 내용을 확인했습니다. (체크 상자)")
    page_no(1); c.showPage()

    c.setFont(F, 18); c.drawString(60, H - 90, "2쪽 — 표와 목록")
    rows = [("항목", "형식", "동작"),
            ("문서", "PDF", "서명·필기·저장"),
            ("표계산", "XLSX", "셀 편집·수식·차트"),
            ("발표", "PPTX", "미리보기"),
            ("코드", "PY", "실행·자동채점")]
    ty = H - 140
    for r in rows:
        c.setFont(F, 11)
        for j, cell in enumerate(r):
            c.drawString(70 + j * 140, ty, cell)
        c.line(60, ty - 6, 470, ty - 6)
        ty -= 26
    c.setFont(F, 11)
    c.drawString(60, ty - 30, "긴 문서에서 페이지 이동이 잘 되는지 확인하려고 3쪽까지 만들었습니다.")
    page_no(2); c.showPage()

    c.setFont(F, 18); c.drawString(60, H - 90, "3쪽 — 찾기 확인용")
    c.setFont(F, 11)
    c.drawString(60, H - 140, "아래 낱말을 위쪽 찾기 상자에 넣어 보세요.")
    c.setFont(F, 16); c.drawString(60, H - 180, "검색테스트")
    c.setFont(F, 11)
    for i in range(12):
        c.drawString(60, H - 220 - i * 18, "보통 문단 %02d — 가나다라마바사 아자차카타파하 ABC abc 0123456789" % (i + 1))
    page_no(3); c.showPage()
    c.save()
    done("문서-샘플.pdf")


# ---------------------------------------------------------------- DOCX / 구형 DOC
def make_docx():
    def para(text, style=None, bold=False, size=None):
        rpr = ""
        if bold: rpr += "<w:b/>"
        if size: rpr += '<w:sz w:val="%d"/>' % (size * 2)
        ppr = '<w:pPr><w:pStyle w:val="%s"/></w:pPr>' % style if style else ""
        return ('<w:p>%s<w:r>%s<w:t xml:space="preserve">%s</w:t></w:r></w:p>'
                % (ppr, "<w:rPr>%s</w:rPr>" % rpr if rpr else "", text))

    def cell(text, bold=False):
        return ('<w:tc><w:tcPr><w:tcW w:w="2600" w:type="dxa"/></w:tcPr>%s</w:tc>'
                % para(text, bold=bold))

    table_rows = [["항목", "확장자", "동작"],
                  ["문서", ".docx", "미리보기"],
                  ["표계산", ".xlsx", "편집·수식·차트"],
                  ["발표", ".pptx", "미리보기"]]
    tbl = ('<w:tbl><w:tblPr><w:tblW w:w="7800" w:type="dxa"/>'
           '<w:tblBorders>' + "".join(
               '<w:%s w:val="single" w:sz="6" w:color="999999"/>' % s
               for s in ("top", "left", "bottom", "right", "insideH", "insideV")) +
           '</w:tblBorders></w:tblPr>' +
           "".join("<w:tr>%s</w:tr>" % "".join(cell(t, bold=(i == 0)) for t in row)
                   for i, row in enumerate(table_rows)) +
           "</w:tbl>")

    body = (para("만능파일교실 Word 테스트 문서", "Heading1")
            + para("이 파일은 .docx 미리보기를 확인하려고 만든 샘플입니다.")
            + para("굵은 글씨 문단입니다.", bold=True)
            + para("글자 크기를 키운 문단입니다.", size=18)
            + para("표 미리보기", "Heading2")
            + tbl
            + para("")
            + para("긴 문단 — " + "가나다라마바사 아자차카타파하 ABC abc 0123456789. " * 8)
            + para("문서 끝. 통합 검색으로 '워드검색테스트' 를 찾아보세요.")
            + para("워드검색테스트"))

    document = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
                '<w:body>%s<w:sectPr><w:pgSz w:w="11906" w:h="16838"/>'
                '<w:pgMar w:top="1418" w:right="1134" w:bottom="1418" w:left="1134"/>'
                '</w:sectPr></w:body></w:document>' % body)

    styles = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
              '<w:styles xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
              '<w:style w:type="paragraph" w:styleId="Heading1"><w:name w:val="heading 1"/>'
              '<w:pPr><w:outlineLvl w:val="0"/></w:pPr><w:rPr><w:b/><w:sz w:val="40"/></w:rPr></w:style>'
              '<w:style w:type="paragraph" w:styleId="Heading2"><w:name w:val="heading 2"/>'
              '<w:pPr><w:outlineLvl w:val="1"/></w:pPr><w:rPr><w:b/><w:sz w:val="30"/></w:rPr></w:style>'
              '</w:styles>')

    ct = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
          '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
          '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
          '<Default Extension="xml" ContentType="application/xml"/>'
          '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
          '<Override PartName="/word/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.styles+xml"/>'
          '</Types>')
    rels = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
            '</Relationships>')
    drels = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
             '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
             '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>'
             '</Relationships>')

    with zipfile.ZipFile(p("문서-샘플.docx"), "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", ct)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", document)
        z.writestr("word/styles.xml", styles)
        z.writestr("word/_rels/document.xml.rels", drels)
    done("문서-샘플.docx")


def make_doc_legacy():
    """구형 .doc 은 단위 테스트 픽스처(build-doc.js)를 node 로 그대로 쓴다 — 규격을 두 벌 두지 않는다."""
    script = os.path.join(REPO, "tests", "fixtures", "build-doc.js").replace("\\", "/")
    target = p("문서-구형샘플.doc").replace("\\", "/")
    text = ("만능파일교실 구형 워드(.doc) 테스트 문서입니다. 이 뷰어는 글자만 뽑아 보여 줍니다"
            "(표·그림·서식 제외). 통합 검색으로 구형워드검색테스트 를 찾아보세요.")
    js = ("const {buildWordDoc}=require(%s);require('fs').writeFileSync(%s,Buffer.from(buildWordDoc({text:%s})));"
          % (json.dumps(script), json.dumps(target), json.dumps(text)))
    subprocess.run(["node", "-e", js], check=True, shell=(os.name == "nt"))
    done("문서-구형샘플.doc")


# ---------------------------------------------------------------- XLSX / CSV
def make_xlsx():
    import random
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    random.seed(7)
    names = ["김하늘", "이바다", "박구름", "최나무", "정바람", "강별", "윤햇살", "임소리", "한겨울", "오여름"]

    wb = Workbook()
    ws = wb.active; ws.title = "성적표"
    ws.append(["번호", "이름", "국어", "영어", "수학", "합계", "평균"])
    for i, n in enumerate(names, 1):
        k, e, m = (random.randint(55, 100) for _ in range(3))
        r = i + 1
        ws.append([i, n, k, e, m, "=SUM(C%d:E%d)" % (r, r), "=ROUND(AVERAGE(C%d:E%d),1)" % (r, r)])
    last = len(names) + 1
    ws.append([])
    ws.cell(row=last + 2, column=2, value="과목 평균")
    for col in range(3, 6):
        L = get_column_letter(col)
        ws.cell(row=last + 2, column=col, value="=ROUND(AVERAGE(%s2:%s%d),1)" % (L, L, last))
    thin = Side(style="thin", color="BBBBBB")
    for c in ws[1]:
        c.font = Font(bold=True, color="FFFFFF")
        c.fill = PatternFill("solid", fgColor="4472C4")
        c.alignment = Alignment(horizontal="center")
    for row in ws.iter_rows(min_row=1, max_row=last, max_col=7):
        for c in row:
            c.border = Border(top=thin, bottom=thin, left=thin, right=thin)
    for i, w in enumerate([6, 12, 8, 8, 8, 8, 8], 1):
        ws.column_dimensions[get_column_letter(i)].width = w
    ws.freeze_panes = "A2"

    ws2 = wb.create_sheet("월별판매")
    ws2.append(["월", "노트북", "태블릿", "휴대폰"])
    for m in range(1, 13):
        ws2.append(["%d월" % m, random.randint(20, 90), random.randint(10, 60), random.randint(40, 130)])
    ws2.column_dimensions["A"].width = 8

    ws3 = wb.create_sheet("메모")
    ws3["A1"] = "시트 전환·검색·범위 통계·내보내기를 확인해 보세요."
    ws3["A3"] = "셀 편집, 자동 채우기, 필터, 정렬, 조건부 서식, 차트, 미니 피벗도 여기서 실험하면 됩니다."
    ws3["A5"] = "엑셀검색테스트"
    wb.save(p("표계산-샘플.xlsx"))
    done("표계산-샘플.xlsx")

    import datetime
    rows = ["날짜,지점,상품,수량,단가,금액,담당자,비고"]
    base = datetime.date(2026, 1, 5)
    shops = ["서울점", "부산점", "대전점", "광주점"]
    items = ["노트북", "태블릿", "휴대폰", "이어폰", "충전기"]
    for i in range(40):
        d = base + datetime.timedelta(days=i * 3)
        q = random.randint(1, 12); price = random.choice([12000, 45000, 890000, 1290000, 330000])
        rows.append("%s,%s,%s,%d,%d,%d,%s,%s" % (
            d.isoformat(), shops[i % 4], items[i % 5], q, price, q * price,
            names[i % len(names)], "정상" if i % 7 else "반품"))
    write_text("표계산-샘플.csv", "\n".join(rows) + "\n", bom=True)
    done("표계산-샘플.csv")


# ---------------------------------------------------------------- PPTX
def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "만능파일교실 발표 자료 샘플"
    s.placeholders[1].text = "PPTX 미리보기 확인용 · 2026"

    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "확인할 내용"
    tf = s.placeholders[1].text_frame
    tf.text = "슬라이드 넘기기"
    for t in ["글머리 기호와 줄바꿈", "표가 들어간 슬라이드", "도형과 색",
              "EXE + 파워포인트 설치 시 PDF 정확 변환"]:
        tf.add_paragraph().text = t

    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "표 슬라이드"
    data = [["항목", "확장자", "동작"], ["문서", ".docx", "미리보기"],
            ["표계산", ".xlsx", "편집·저장"], ["한글", ".hwpx", "간이 미리보기"]]
    tbl = s.shapes.add_table(len(data), 3, Inches(0.8), Inches(1.8), Inches(8), Inches(2.5)).table
    for r in range(len(data)):
        for c in range(3):
            cell = tbl.cell(r, c); cell.text = data[r][c]
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(14)

    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "도형 슬라이드"
    for i, (shape, x) in enumerate([(MSO_SHAPE.ROUNDED_RECTANGLE, 0.8),
                                    (MSO_SHAPE.OVAL, 3.6), (MSO_SHAPE.RIGHT_ARROW, 6.4)]):
        sp = s.shapes.add_shape(shape, Inches(x), Inches(2.2), Inches(2.4), Inches(1.6))
        sp.text_frame.text = ["시작", "처리", "끝"][i]
    prs.save(p("발표-샘플.pptx"))
    done("발표-샘플.pptx")


# ---------------------------------------------------------------- HWPX
def make_hwpx():
    """OWPML(zip+XML) 최소 구조 — office-doc-viewers.js 의 renderHwpx 가 읽는 부분만 채운다."""
    NS = ('xmlns:hs="http://www.hancom.co.kr/hwpml/2011/section" '
          'xmlns:hp="http://www.hancom.co.kr/hwpml/2011/paragraph"')

    def run(text, cid="1"):
        return '<hp:run charPrIDRef="%s"><hp:t>%s</hp:t></hp:run>' % (cid, text)

    def para(text, cid="1", pid="0"):
        return '<hp:p paraPrIDRef="%s">%s</hp:p>' % (pid, run(text, cid))

    def tcell(text):
        return ('<hp:tc><hp:cellSpan colSpan="1" rowSpan="1"/><hp:subList>%s</hp:subList></hp:tc>'
                % para(text))

    rows = [["항목", "확장자", "동작"], ["한글", ".hwp", "hwp.js 뷰어"],
            ["한글(신형)", ".hwpx", "간이 미리보기"]]
    tbl = "<hp:tbl>" + "".join("<hp:tr>%s</hp:tr>" % "".join(tcell(c) for c in r) for r in rows) + "</hp:tbl>"

    section = ('<?xml version="1.0" encoding="UTF-8"?><hs:sec %s>' % NS
               + para("만능파일교실 HWPX 테스트 문서", cid="2", pid="1")
               + para("")
               + para("이 파일은 .hwpx 간이 미리보기를 확인하려고 만든 샘플입니다.")
               + para("굵은 글씨 · 색 글씨도 확인해 보세요.", cid="3")
               + para("")
               + '<hp:p paraPrIDRef="0"><hp:run charPrIDRef="1">%s</hp:run></hp:p>' % tbl
               + para("")
               + para("긴 문단 — " + "가나다라마바사 아자차카타파하 ABC abc 0123456789. " * 5)
               + para("한글검색테스트")
               + "</hs:sec>")

    header = ('<?xml version="1.0" encoding="UTF-8"?>'
              '<hh:head xmlns:hh="http://www.hancom.co.kr/hwpml/2011/head">'
              '<hh:refList><hh:charProperties>'
              '<hh:charPr id="1" height="1000" textColor="#222222"/>'
              '<hh:charPr id="2" height="1800" textColor="#1a4fa0"><hh:bold/></hh:charPr>'
              '<hh:charPr id="3" height="1100" textColor="#c0392b"><hh:bold/><hh:underline type="SOLID"/></hh:charPr>'
              '</hh:charProperties><hh:paraProperties>'
              '<hh:paraPr id="0"><hh:align horizontal="LEFT"/></hh:paraPr>'
              '<hh:paraPr id="1"><hh:align horizontal="CENTER"/></hh:paraPr>'
              '</hh:paraProperties></hh:refList></hh:head>')

    hpf = ('<?xml version="1.0" encoding="UTF-8"?>'
           '<opf:package xmlns:opf="http://www.idpf.org/2007/opf/"><opf:manifest>'
           '<opf:item id="header" href="Contents/header.xml" media-type="application/xml"/>'
           '<opf:item id="section0" href="Contents/section0.xml" media-type="application/xml"/>'
           '</opf:manifest></opf:package>')

    manifest = ('<?xml version="1.0" encoding="UTF-8"?>'
                '<manifest xmlns="http://www.hancom.co.kr/schema/2011/hpf">'
                '<file-entry full-path="/" media-type="application/hwp+zip"/></manifest>')

    with zipfile.ZipFile(p("한글-샘플.hwpx"), "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr(zipfile.ZipInfo("mimetype"), "application/hwp+zip", zipfile.ZIP_STORED)
        z.writestr("version.xml", '<?xml version="1.0" encoding="UTF-8"?><hv:HCFVersion '
                   'xmlns:hv="http://www.hancom.co.kr/hwpml/2011/version" tagetApplication="WORDPROCESSOR" '
                   'major="5" minor="1" micro="0" buildNumber="0"/>')
        z.writestr("META-INF/manifest.xml", manifest)
        z.writestr("META-INF/container.rdf", '<?xml version="1.0" encoding="UTF-8"?>'
                   '<rdf:RDF xmlns:rdf="http://www.w3.org/1999/02/22-rdf-syntax-ns#"/>')
        z.writestr("Contents/content.hpf", hpf)
        z.writestr("Contents/header.xml", header)
        z.writestr("Contents/section0.xml", section)
        # 본문 파싱이 실패해도 텍스트 미리보기로 폴백되는지 함께 확인할 수 있게 넣어 둔다.
        z.writestr("Preview/PrvText.txt",
                   "만능파일교실 HWPX 테스트 문서\n간이 미리보기가 실패하면 이 텍스트가 보입니다.\n".encode("utf-8"))
    done("한글-샘플.hwpx")


# ---------------------------------------------------------------- 이미지
def kfont(size, bold=False):
    from PIL import ImageFont
    try:
        return ImageFont.truetype(KFONTB if bold else KFONT, size)
    except Exception:
        return ImageFont.load_default()


def poster(w, h, title, sub, hue=(30, 90, 190)):
    """어느 형식으로 저장해도 알아보기 쉬운 그림 한 장(하늘 그러데이션 + 산 + 제목)."""
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (w, h))
    d = ImageDraw.Draw(img)
    for y in range(h):
        t = y / max(1, h - 1)
        d.line([(0, y), (w, y)], fill=(int(hue[0] + (250 - hue[0]) * t),
                                       int(hue[1] + (240 - hue[1]) * t),
                                       int(hue[2] + (230 - hue[2]) * t)))
    d.ellipse([w * .55, h * .12, w * .95, h * .52], fill=(255, 214, 102))
    d.polygon([(0, h), (w * .35, h * .45), (w * .6, h)], fill=(46, 125, 90))
    d.polygon([(w * .3, h), (w * .68, h * .35), (w, h)], fill=(30, 95, 70))
    d.text((28, 24), title, font=kfont(max(18, w // 16), True), fill=(255, 255, 255))
    d.text((28, 24 + max(18, w // 16) + 12), sub, font=kfont(max(12, w // 30)), fill=(240, 240, 240))
    return img


def make_images():
    from PIL import Image, ImageDraw
    poster(1280, 800, "이미지 테스트 (PNG)", "회전·자르기·모자이크·보정을 확인해 보세요").save(p("이미지-풍경.png"))
    done("이미지-풍경.png")
    poster(1600, 1067, "이미지 테스트 (JPG)", "사진처럼 큰 이미지 · 확대/축소 확인",
           (150, 60, 40)).save(p("이미지-사진.jpg"), quality=88)
    done("이미지-사진.jpg")
    poster(800, 500, "WEBP 테스트", "웹 이미지 형식", (60, 40, 140)).save(p("이미지-웹.webp"))
    done("이미지-웹.webp")
    poster(640, 400, "BMP 테스트", "무압축 비트맵", (20, 110, 110)).save(p("이미지-비트맵.bmp"))
    done("이미지-비트맵.bmp")
    poster(256, 256, "ICO", "아이콘", (200, 80, 30)).save(
        p("이미지-아이콘.ico"), sizes=[(16, 16), (32, 32), (64, 64), (128, 128), (256, 256)])
    done("이미지-아이콘.ico")

    img = Image.new("RGBA", (600, 400), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    d.ellipse([50, 50, 350, 350], fill=(255, 99, 132, 210))
    d.rectangle([250, 150, 550, 350], fill=(54, 162, 235, 190))
    d.text((60, 20), "투명 배경 PNG", font=kfont(28, True), fill=(20, 20, 20, 255))
    img.save(p("이미지-투명배경.png"))
    done("이미지-투명배경.png")

    frames = []
    for i in range(12):
        f = Image.new("RGB", (400, 300), (250, 250, 252))
        dd = ImageDraw.Draw(f)
        x = 30 + i * 28
        dd.ellipse([x, 120 + int(40 * math.sin(i / 2)), x + 60, 180 + int(40 * math.sin(i / 2))],
                   fill=(230, 80, 90))
        dd.text((20, 20), "움직이는 GIF %02d" % (i + 1), font=kfont(22, True), fill=(40, 40, 60))
        frames.append(f)
    frames[0].save(p("이미지-움직임.gif"), save_all=True, append_images=frames[1:], duration=110, loop=0)
    done("이미지-움직임.gif")

    try:                                        # AVIF 는 Pillow 빌드에 따라 없을 수 있다
        poster(800, 500, "AVIF 테스트", "최신 이미지 형식", (90, 30, 120)).save(p("이미지-avif.avif"))
        done("이미지-avif.avif")
    except Exception as e:
        skipped.append("이미지-avif.avif: %s" % e)
        print("  - 건너뜀: AVIF :", e)

    write_text("이미지-벡터.svg", '''<svg xmlns="http://www.w3.org/2000/svg" width="640" height="400" viewBox="0 0 640 400">
  <defs><linearGradient id="g" x1="0" y1="0" x2="0" y2="1">
    <stop offset="0" stop-color="#1e5fbe"/><stop offset="1" stop-color="#8ec5ff"/>
  </linearGradient></defs>
  <rect width="640" height="400" fill="url(#g)"/>
  <circle cx="500" cy="110" r="60" fill="#ffd666"/>
  <polygon points="0,400 220,180 380,400" fill="#2e7d5a"/>
  <polygon points="200,400 420,150 640,400" fill="#1f6049"/>
  <text x="30" y="60" font-size="34" font-family="Malgun Gothic, sans-serif" fill="#ffffff">벡터 이미지 (SVG)</text>
  <text x="30" y="96" font-size="18" font-family="Malgun Gothic, sans-serif" fill="#eaf2ff">확대해도 깨지지 않는지 확인해 보세요</text>
</svg>
''')
    done("이미지-벡터.svg")

    for i in range(1, 9):                       # 폴더 안 이미지 격자 모아보기 확인용
        poster(600, 400, "사진 %02d" % i, "폴더 격자 모아보기 확인",
               (30 + i * 20, 60 + i * 15, 190 - i * 12)).save(p("사진모음", "사진-%02d.jpg" % i), quality=85)
    done("사진모음/ (사진 8장)")


# ---------------------------------------------------------------- 영상 / 자막 / 오디오
def make_video():
    import cv2
    import numpy as np
    from PIL import Image, ImageDraw

    def frames(n, w, h, label):
        for i in range(n):
            img = Image.new("RGB", (w, h), (18, 22, 34))
            d = ImageDraw.Draw(img)
            t = i / n
            d.rectangle([0, 0, w, int(h * .28)], fill=(30, 90, 190))
            d.text((24, 18), label, font=kfont(int(h / 12), True), fill=(255, 255, 255))
            d.text((24, int(h * .34)), "%05.2f 초" % (i / 24), font=kfont(int(h / 10), True), fill=(250, 220, 90))
            cx = int(w * .15 + (w * .7) * abs(math.sin(math.pi * t * 2)))
            d.ellipse([cx - 30, int(h * .72) - 30, cx + 30, int(h * .72) + 30], fill=(235, 80, 100))
            d.rectangle([0, h - 14, int(w * t), h], fill=(120, 200, 140))
            yield np.array(img)[:, :, ::-1].copy()          # PIL(RGB) → OpenCV(BGR)

    def write(name, fourcc, w, h, n, label):
        # OpenCV 의 FFmpeg 백엔드는 한글이 든 경로에서 인코더 초기화가 실패할 수 있어
        # 임시 폴더(ASCII 경로)에 만든 뒤 옮긴다.
        ext = os.path.splitext(name)[1]
        tmp = os.path.join(tempfile.gettempdir(), "mn-sample-%d%s" % (int(time.time() * 1000), ext))
        vw = cv2.VideoWriter(tmp, cv2.VideoWriter_fourcc(*fourcc), 24, (w, h))
        for f in frames(n, w, h, label):
            vw.write(f)
        vw.release()
        if not os.path.exists(tmp) or os.path.getsize(tmp) < 1024:
            raise RuntimeError("%s 인코딩 실패(코덱 %s 없음)" % (name, fourcc))
        shutil.move(tmp, p(name))
        done(name)

    write("영상-샘플.mp4", "avc1", 640, 360, 24 * 8, "영상 재생 테스트 (MP4/H.264)")
    write("영상-웹형식.webm", "VP80", 640, 360, 24 * 6, "영상 재생 테스트 (WebM/VP8)")
    write("영상-변환필요.avi", "MJPG", 480, 270, 24 * 4, "브라우저가 못 여는 AVI")


def make_subtitles():
    write_text("영상-샘플.srt", """1
00:00:00,500 --> 00:00:02,500
안녕하세요! 자막 테스트입니다.

2
00:00:02,600 --> 00:00:04,800
영상과 이름이 같으면 자동으로 연결됩니다.

3
00:00:05,000 --> 00:00:07,000
자막 글자 크기도 바꿔 보세요.

4
00:00:07,100 --> 00:00:08,000
끝!
""")
    done("영상-샘플.srt (mp4 와 자동 연결)")

    write_text("영상-웹형식.vtt", """WEBVTT

00:00:00.500 --> 00:00:02.500
WebVTT 자막입니다.

00:00:02.600 --> 00:00:05.000
수동 연결 / 표시·숨기기를 확인해 보세요.
""")
    done("영상-웹형식.vtt (webm 과 자동 연결)")

    write_text("자막-구형.smi", """<SAMI>
<HEAD><TITLE>SMI 자막 테스트</TITLE>
<STYLE TYPE="text/css">
<!--
P { font-family:Malgun Gothic; font-size:20pt; color:white; text-align:center; }
.KRCC { Name:한국어; lang:ko-KR; SAMIType:CC; }
-->
</STYLE></HEAD>
<BODY>
<SYNC Start=500><P Class=KRCC>구형 SMI 자막입니다.
<SYNC Start=2600><P Class=KRCC>수동으로 연결해서 확인해 보세요.
<SYNC Start=5000><P Class=KRCC>&nbsp;
</BODY>
</SAMI>
""")
    done("자막-구형.smi")


def make_audio():
    sr, dur = 44100, 4.0                        # 도·미·솔·도 아르페지오(실제 소리)
    notes = [261.63, 329.63, 392.00, 523.25]
    with wave.open(p("오디오-샘플.wav"), "w") as w:
        w.setnchannels(2); w.setsampwidth(2); w.setframerate(sr)
        buf = bytearray()
        for i in range(int(sr * dur)):
            t = i / sr
            f = notes[min(len(notes) - 1, int(t / (dur / len(notes))))]
            env = min(1.0, (t % (dur / len(notes))) * 8) * max(0.0, 1 - (t / dur) * 0.6)
            v = int(12000 * env * math.sin(2 * math.pi * f * t))
            buf += struct.pack("<hh", v, int(v * 0.7))
        w.writeframes(bytes(buf))
    done("오디오-샘플.wav")

    # MP3 인코더가 없으므로 무음 프레임을 직접 쌓는다(재생·탐색 UI 확인용, 소리는 없음).
    # 0xFF 0xFB = MPEG1 Layer3, 0x90 = 128kbps/44.1kHz, 0x64 = 조인트 스테레오 → 프레임 417바이트.
    frame = bytes([0xFF, 0xFB, 0x90, 0x64]) + b"\x00" * 413
    with open(p("오디오-무음.mp3"), "wb") as f:
        f.write(frame * int(4.0 * 44100 / 1152))
    done("오디오-무음.mp3 (무음)")


# ---------------------------------------------------------------- 파이썬 / 노트북
def make_python():
    write_text("파이썬-계산기.py", '''"""표준입력·반복문·함수를 함께 확인하는 실행 테스트용 코드."""

def 평균(값들):
    return sum(값들) / len(값들) if 값들 else 0


def main():
    print("점수를 공백으로 구분해 입력하세요. (예: 90 85 77)")
    줄 = input("점수: ").strip()
    점수 = [int(x) for x in 줄.split()] if 줄 else [90, 85, 77]
    print("입력한 점수:", 점수)
    print("합계 :", sum(점수))
    print("평균 : %.1f" % 평균(점수))
    for i, s in enumerate(sorted(점수, reverse=True), 1):
        print("%d등 → %d점" % (i, s))


if __name__ == "__main__":
    main()
''')
    done("파이썬-계산기.py")

    write_text("파이썬-그래프.py", '''"""matplotlib 그래프가 실행 결과로 뜨는지 확인하는 코드."""
import matplotlib.pyplot as plt

월 = list(range(1, 13))
판매 = [42, 51, 63, 58, 72, 88, 95, 90, 76, 64, 55, 61]

plt.figure(figsize=(7, 4))
plt.plot(월, 판매, marker="o")
plt.title("Monthly Sales 2026")
plt.xlabel("month")
plt.ylabel("sales")
plt.grid(alpha=.3)
plt.tight_layout()
plt.show()

print("최고 판매:", max(판매), "월:", 월[판매.index(max(판매))])
''')
    done("파이썬-그래프.py")

    write_text("파이썬-오류.py", '''"""오류 표시(빨간 traceback)와 실시간 진단을 확인하는 코드."""

def 나누기(a, b):
    return a / b


print("여기까지는 정상 출력")
print(나누기(10, 0))      # ZeroDivisionError 가 납니다
print("여기는 실행되지 않습니다")
''')
    done("파이썬-오류.py")

    write_text(os.path.join("파이썬프로젝트", "__init__.py"), "")
    write_text(os.path.join("파이썬프로젝트", "도우미.py"), '''"""옆 파일 자동완성·자동 import 확인용 모듈."""


class 학생:
    def __init__(self, 이름, 점수):
        self.이름 = 이름
        self.점수 = 점수

    def 요약(self):
        return f"{self.이름}: {self.점수}점"


def 상위(학생들, n=3):
    return sorted(학생들, key=lambda s: s.점수, reverse=True)[:n]
''')
    write_text(os.path.join("파이썬프로젝트", "main.py"), '''"""편집기에서 '학생' 을 입력하면 옆 파일이 자동완성/자동 import 되는지 확인."""
from 파이썬프로젝트.도우미 import 학생, 상위

명단 = [학생("김하늘", 92), 학생("이바다", 78), 학생("박구름", 85), 학생("최나무", 61)]

for s in 상위(명단):
    print(s.요약())
''')
    done("파이썬프로젝트/ (도우미.py · main.py)")


def make_notebook():
    def code(src):
        return {"cell_type": "code", "execution_count": None, "metadata": {},
                "outputs": [], "source": src.splitlines(True)}

    def md(src):
        return {"cell_type": "markdown", "metadata": {}, "source": src.splitlines(True)}

    nb = {
        "cells": [
            md("# 주피터 노트북 테스트\n\n셀 편집·실행·순서 바꾸기를 확인하는 샘플입니다.\n\n"
               "- 코드 셀과 마크다운 셀이 섞여 있습니다\n- 셀별 실행과 커널 상태 유지를 확인해 보세요"),
            code("이름 = '만능파일교실'\n합계 = sum(range(1, 101))\nprint(이름, '합계:', 합계)\n"),
            md("## 커널 상태 유지 확인\n\n위 셀에서 만든 변수를 아래 셀에서 그대로 씁니다."),
            code("print(이름[::-1])\nprint('합계의 두 배 =', 합계 * 2)\n"),
            md("## 그래프 출력"),
            code("import matplotlib.pyplot as plt\n\nxs = list(range(0, 20))\nys = [x * x for x in xs]\n"
                 "plt.plot(xs, ys, marker='o')\nplt.title('y = x^2')\nplt.grid(alpha=.3)\nplt.show()\n"),
            md("## 표 출력 (pandas 가 있을 때)"),
            code("try:\n    import pandas as pd\n    df = pd.DataFrame({'이름': ['하늘', '바다', '구름'], "
                 "'점수': [92, 78, 85]})\n    display(df)\nexcept Exception as e:\n    print('pandas 없음:', e)\n"),
            md("### 메모\n\n셀 메모 보관·재사용, 드래그로 순서 바꾸기도 여기서 실험해 보세요."),
        ],
        "metadata": {"kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"},
                     "language_info": {"name": "python", "version": "3.13"}},
        "nbformat": 4, "nbformat_minor": 5,
    }
    write_text("노트북-샘플.ipynb", json.dumps(nb, ensure_ascii=False, indent=1))
    done("노트북-샘플.ipynb")


# ---------------------------------------------------------------- SQLite
def make_db():
    import random
    path = p("데이터베이스-샘플.db")
    if os.path.exists(path):
        os.remove(path)
    random.seed(11)
    con = sqlite3.connect(path)
    c = con.cursor()
    c.execute("CREATE TABLE 학생(id INTEGER PRIMARY KEY, 이름 TEXT, 학년 INTEGER, 반 INTEGER)")
    c.execute("CREATE TABLE 성적(id INTEGER PRIMARY KEY, 학생id INTEGER, 과목 TEXT, 점수 INTEGER)")
    c.execute("CREATE TABLE 과목(코드 TEXT PRIMARY KEY, 이름 TEXT, 시수 INTEGER)")
    names = ["김하늘", "이바다", "박구름", "최나무", "정바람", "강별", "윤햇살", "임소리"]
    for i, n in enumerate(names, 1):
        c.execute("INSERT INTO 학생 VALUES(?,?,?,?)", (i, n, 1 + i % 3, 1 + i % 4))
    subjects = [("KOR", "국어", 4), ("ENG", "영어", 4), ("MAT", "수학", 5), ("SCI", "과학", 3)]
    c.executemany("INSERT INTO 과목 VALUES(?,?,?)", subjects)
    rid = 1
    for sid in range(1, len(names) + 1):
        for _, sub, _n in subjects:
            c.execute("INSERT INTO 성적 VALUES(?,?,?,?)", (rid, sid, sub, random.randint(50, 100)))
            rid += 1
    c.execute("CREATE VIEW 학생평균 AS SELECT 학생.이름, ROUND(AVG(성적.점수),1) AS 평균 "
              "FROM 학생 JOIN 성적 ON 학생.id=성적.학생id GROUP BY 학생.id")
    con.commit(); con.close()
    done("데이터베이스-샘플.db")


# ---------------------------------------------------------------- 텍스트 / 코드 / 설정
def make_texts():
    write_text("마크다운-샘플.md", '''# 마크다운 테스트 문서

만능파일교실의 **마크다운 렌더**를 확인하는 샘플입니다.

## 목록

- 첫째 항목
- 둘째 항목
  - 들여쓴 항목
1. 번호 목록
2. 두 번째

## 표

| 항목 | 확장자 | 동작 |
|------|--------|------|
| 문서 | `.pdf` | 서명·필기·저장 |
| 표계산 | `.xlsx` | 편집·수식·차트 |
| 코드 | `.py` | 실행·자동채점 |

## 코드 블록

```python
def 인사(이름):
    return f"안녕하세요, {이름}님!"

print(인사("선생님"))
```

## 인용과 강조

> 인용문입니다. 이 줄이 들여쓰기 되어 보이면 정상입니다.

*기울임*, **굵게**, ~~취소선~~, `인라인 코드`.

---

[링크 테스트](https://example.com) · 마크다운검색테스트
''')
    done("마크다운-샘플.md")

    write_text("웹페이지-샘플.html", '''<!doctype html>
<html lang="ko">
<head>
<meta charset="utf-8">
<title>HTML 미리보기 테스트</title>
<style>
  body { font-family: "Malgun Gothic", system-ui, sans-serif; margin: 32px; color:#222; }
  h1 { color:#1e5fbe; }
  table { border-collapse: collapse; margin-top: 12px; }
  td, th { border:1px solid #bbb; padding:6px 12px; }
  .box { background:#eef4ff; padding:16px; border-radius:10px; margin-top:16px; }
  button { padding:8px 14px; border-radius:8px; border:1px solid #1e5fbe; background:#fff; cursor:pointer; }
</style>
</head>
<body>
  <h1>HTML 미리보기 테스트</h1>
  <p>격리된 미리보기로 열리는지, 스타일이 적용되는지 확인해 보세요.</p>
  <div class="box">
    <button onclick="document.getElementById('out').textContent = '버튼이 눌렸습니다 · ' + new Date().toLocaleTimeString()">눌러보기</button>
    <p id="out">스크립트 동작 확인용 영역</p>
  </div>
  <table>
    <tr><th>항목</th><th>값</th></tr>
    <tr><td>글자</td><td>가나다라 ABC 123</td></tr>
    <tr><td>이모지</td><td>📄 📊 🎬</td></tr>
  </table>
  <p>HTML검색테스트</p>
</body>
</html>
''')
    done("웹페이지-샘플.html")

    write_text("텍스트-샘플.txt", "\n".join(
        ["만능파일교실 텍스트 파일 테스트", "=" * 34, "",
         "이 파일은 일반 텍스트 보기를 확인하는 샘플입니다.",
         "한글 · English · 123456 · !@#$%^&*() · 특수문자 ①②③ ㉠㉡ ★☆ ",
         "탭\t문자와\t줄맞춤도 확인해 보세요.", ""]
        + ["%03d 번째 줄 — 통합 검색과 줄 번호 이동(Ctrl+G)을 확인해 보세요." % i for i in range(1, 121)]
        + ["", "텍스트검색테스트", ""]))
    done("텍스트-샘플.txt")

    write_text("코드-샘플.js", '''"use strict";

/* 자바스크립트 구문 강조 · 줄 번호 확인용 샘플 */

const 과목 = ["국어", "영어", "수학"];

class 성적표 {
  constructor(이름) {
    this.이름 = 이름;
    this.점수 = new Map();
  }
  넣기(과목명, 점수) {
    this.점수.set(과목명, 점수);
    return this;
  }
  get 평균() {
    const 값 = [...this.점수.values()];
    return 값.length ? 값.reduce((a, b) => a + b, 0) / 값.length : 0;
  }
}

async function 불러오기(url) {
  try {
    const res = await fetch(url);
    if (!res.ok) throw new Error(`HTTP ${res.status}`);
    return await res.json();
  } catch (err) {
    console.warn("실패:", err.message);
    return null;
  }
}

const 표 = new 성적표("김하늘");
과목.forEach((s, i) => 표.넣기(s, 80 + i * 5));
console.log(표.이름, 표.평균.toFixed(1));
''')
    done("코드-샘플.js")

    write_text("코드-샘플.ts", '''/* 타입스크립트 구문 강조 확인용 샘플 */

export type 등급 = "A" | "B" | "C" | "D";

export interface 학생 {
  id: number;
  이름: string;
  점수: number[];
  메모?: string;
}

export function 평균(학생: 학생): number {
  if (!학생.점수.length) return 0;
  return 학생.점수.reduce((a, b) => a + b, 0) / 학생.점수.length;
}

export function 등급매기기(점수: number): 등급 {
  if (점수 >= 90) return "A";
  if (점수 >= 80) return "B";
  if (점수 >= 70) return "C";
  return "D";
}

const 명단: 학생[] = [
  { id: 1, 이름: "김하늘", 점수: [92, 88, 95] },
  { id: 2, 이름: "이바다", 점수: [70, 75, 68], 메모: "보충 필요" },
];

명단.forEach((s) => console.log(s.이름, 등급매기기(평균(s))));
''')
    done("코드-샘플.ts")

    write_text("설정-샘플.json", json.dumps({
        "이름": "만능파일교실 테스트 설정",
        "버전": "1.0.0",
        "사용": True,
        "허용확장자": ["pdf", "docx", "xlsx", "pptx", "hwpx", "py", "ipynb"],
        "화면": {"테마": "자동", "글자크기": 15, "사이드바": {"열림": True, "너비": 260}},
        "서버": {"주소": "127.0.0.1", "포트": 17645, "대체포트": [18645, 19645, 27645]},
        "메모": "JSON 구문 강조와 접기를 확인해 보세요. json검색테스트",
    }, ensure_ascii=False, indent=2))
    done("설정-샘플.json")

    write_text("스타일-샘플.css", ''':root {
  --바탕: #ffffff;
  --글자: #1f2430;
  --강조: #1e5fbe;
}

@media (prefers-color-scheme: dark) {
  :root { --바탕: #16181f; --글자: #e8eaf0; }
}

body {
  margin: 0;
  background: var(--바탕);
  color: var(--글자);
  font-family: "Malgun Gothic", system-ui, sans-serif;
}

.카드 {
  display: grid;
  grid-template-columns: repeat(auto-fill, minmax(220px, 1fr));
  gap: 16px;
  padding: 24px;
}

.카드 > article {
  border: 1px solid color-mix(in srgb, var(--글자) 15%, transparent);
  border-radius: 12px;
  padding: 16px;
  transition: transform .18s ease;
}

.카드 > article:hover { transform: translateY(-2px); }
''')
    done("스타일-샘플.css")

    write_text("쿼리-샘플.sql", '''-- SQL 구문 강조 확인용 샘플 (데이터베이스-샘플.db 와 짝)

CREATE TABLE IF NOT EXISTS 학생 (
  id    INTEGER PRIMARY KEY,
  이름  TEXT    NOT NULL,
  학년  INTEGER DEFAULT 1,
  반    INTEGER
);

INSERT INTO 학생 (이름, 학년, 반) VALUES
  ('김하늘', 1, 1),
  ('이바다', 2, 3);

SELECT  학생.이름,
        ROUND(AVG(성적.점수), 1) AS 평균,
        COUNT(*)                 AS 과목수
FROM    학생
JOIN    성적 ON 성적.학생id = 학생.id
WHERE   학생.학년 = 1
GROUP BY 학생.id
HAVING  평균 >= 70
ORDER BY 평균 DESC
LIMIT 10;
''')
    done("쿼리-샘플.sql")

    write_text("데이터-샘플.xml", '''<?xml version="1.0" encoding="UTF-8"?>
<학급 이름="1학년 3반" 담임="정바람">
  <학생 번호="1">
    <이름>김하늘</이름>
    <점수 과목="국어">92</점수>
    <점수 과목="영어">88</점수>
    <메모><![CDATA[특수문자 <>&" 도 그대로 담깁니다]]></메모>
  </학생>
  <학생 번호="2">
    <이름>이바다</이름>
    <점수 과목="국어">70</점수>
    <점수 과목="영어">75</점수>
  </학생>
</학급>
''')
    done("데이터-샘플.xml")

    write_text("설정-샘플.yaml", '''# YAML 구문 강조 확인용 샘플
이름: 만능파일교실
버전: 1.0.0
사용: true

화면:
  테마: 자동
  글자크기: 15
  사이드바:
    열림: true
    너비: 260

허용확장자:
  - pdf
  - docx
  - xlsx
  - pptx
  - py

서버:
  주소: 127.0.0.1
  포트: 17645
  대체포트: [18645, 19645, 27645]

설명: >
  여러 줄 문자열도
  잘 보이는지 확인해 보세요.
''')
    done("설정-샘플.yaml")


# ---------------------------------------------------------------- .lesson / .task / .mnote
def make_lesson():
    """src/js/lesson-replay.js 의 validateLessonPayload 가 받는 최소 스키마(kind:"board")."""
    def stroke(points, color, width, kind="pen"):
        return {"type": kind, "color": color, "width": width,
                "points": [{"x": round(x, 1), "y": round(y, 1)} for x, y in points]}

    keyframes = [{"t": 0, "s": []}]
    t = 300
    keyframes.append({"t": t, "a": {"type": "text", "text": "수업 리플레이 테스트",
                                    "x": 80, "y": 60, "color": "#1f2430", "fontSize": 44}})
    t += 900
    keyframes.append({"t": t, "a": stroke([(80 + i * 12, 120 + math.sin(i / 3) * 2) for i in range(46)],
                                          "#e11d48", 5)})
    t += 900
    keyframes.append({"t": t, "a": stroke([(120 + i * 8, 380 - 120 * math.sin(i / 9)) for i in range(120)],
                                          "#1e5fbe", 4)})
    t += 1600
    keyframes.append({"t": t, "a": stroke([(120 + i * 20, 380) for i in range(48)],
                                          "#ffd400", 26, "highlighter")})
    t += 900
    keyframes.append({"t": t, "a": {"type": "text", "text": "y = sin(x)", "x": 980, "y": 300,
                                    "color": "#1e5fbe", "fontSize": 30}})
    t += 900
    keyframes.append({"t": t, "a": {"type": "rect", "x1": 80, "y1": 480, "x2": 620, "y2": 620,
                                    "color": "#0f9d58", "width": 4}})
    t += 800
    keyframes.append({"t": t, "a": {"type": "arrow", "x1": 640, "y1": 550, "x2": 900, "y2": 550,
                                    "color": "#7c3aed", "width": 5}})
    t += 800
    keyframes.append({"t": t, "a": {"type": "ellipse", "x1": 920, "y1": 470, "x2": 1180, "y2": 630,
                                    "color": "#d97706", "width": 4}})
    t += 900
    keyframes.append({"t": t, "a": {"type": "text", "text": "되감아 재생 · 0.5~4배속을 확인해 보세요",
                                    "x": 120, "y": 530, "color": "#1f2430", "fontSize": 26}})
    t += 1200

    write_text("수업리플레이-샘플.lesson", json.dumps(
        {"format": "manneung-lesson", "version": 1, "kind": "board",
         "createdAt": int(time.time() * 1000), "bg": "#ffffff", "W": 1280, "H": 720,
         "duration": t, "keyframes": keyframes}, ensure_ascii=False))
    done("수업리플레이-샘플.lesson")


def make_task():
    """src/js/task-package.js 의 validateTaskPayload 가 받는 스키마(숨김 테스트·첨부 포함)."""
    data_csv = "이름,점수\n김하늘,92\n이바다,78\n박구름,85\n최나무,61\n정바람,99\n"
    task = {
        "format": "manneung-task", "version": 1, "id": "sample-task-0001",
        "meta": {"title": "평균 구하기", "author": "선생님", "createdAt": "2026-08-06"},
        "problem": {"md": """# 평균 구하기

표준입력으로 **공백으로 구분된 정수**가 한 줄 들어옵니다.
이 수들의 **합계**와 **평균**을 아래 형식으로 출력하세요.

```
합계: 255
평균: 85.0
```

- 평균은 소수점 첫째 자리까지 출력합니다.
- 입력이 비어 있는 경우는 생각하지 않아도 됩니다.

## 참고

같은 폴더의 `data.csv` 파일도 함께 들어 있습니다. `open('data.csv')` 로 열어 볼 수 있어요.

풀고 나면 과제 바의 **✓ 채점** → **📤 제출본 내보내기** 를 눌러 보세요.
"""},
        "starter": {"name": "main.py", "code": '''# 아래 빈칸을 채워 보세요.
수들 = [int(x) for x in input().split()]

합계 = 0          # TODO: 합계를 구하세요
평균 = 0          # TODO: 평균을 구하세요

print("합계:", 합계)
print("평균: %.1f" % 평균)
'''},
        "files": [{"path": "data.csv", "b64": base64.b64encode(data_csv.encode("utf-8")).decode("ascii")}],
        "tests": [
            {"name": "기본 예제", "input": "90 85 80", "expected": "합계: 255\n평균: 85.0"},
            {"name": "한 개", "input": "42", "expected": "합계: 42\n평균: 42.0"},
            {"name": "소수 반올림", "input": "1 2", "expected": "합계: 3\n평균: 1.5"},
            {"name": "숨김 테스트", "input": "7 7 7 7", "expected": "합계: 28\n평균: 7.0", "hidden": True},
        ],
        "options": {},
    }
    write_text("과제-샘플.task", json.dumps(task, ensure_ascii=False, indent=2))
    done("과제-샘플.task")


def make_mnote():
    """src/js/mnote.js 의 mnoteParse 가 받는 스키마(글·표·이미지 블록). 이미지는 파일 안 base64."""
    buf = io.BytesIO()
    poster(560, 320, "블록 문서 안 그림", "이미지 블록 확인용", (120, 40, 90)).save(buf, format="PNG")
    src = "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode("ascii")
    now = int(time.time() * 1000)
    write_text("블록문서-샘플.mnote", json.dumps({
        "format": "manneung-note", "version": 1, "title": "블록 문서 샘플",
        "createdAt": now, "updatedAt": now,
        "blocks": [
            {"id": "text-1", "type": "text",
             "text": "블록 문서(.mnote) 테스트\n\n글·표·이미지 블록을 한 문서에서 섞어 편집할 수 있습니다.\n"
                     "블록을 추가하고 순서를 바꿔 보세요."},
            {"id": "table-1", "type": "table", "header": True,
             "rows": [["항목", "확장자", "동작"],
                      ["문서", ".pdf", "서명·필기·저장"],
                      ["표계산", ".xlsx", "편집·수식·차트"],
                      ["블록 문서", ".mnote", "글·표·이미지 혼합"]]},
            {"id": "image-1", "type": "image", "src": src, "name": "샘플 그림.png",
             "mime": "image/png", "width": "medium",
             "caption": "이미지 블록에는 설명을 달 수 있습니다"},
            {"id": "text-2", "type": "text",
             "text": "HTML·Markdown 내보내기와 되돌리기(Ctrl+Z), 통합 검색도 확인해 보세요.\n\n블록검색테스트"},
        ],
    }, ensure_ascii=False, indent=2))
    done("블록문서-샘플.mnote")


# ---------------------------------------------------------------- 압축
def make_archives():
    with zipfile.ZipFile(p("압축-샘플.zip"), "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("읽어보세요.txt", "압축 파일 안에서도 폴더처럼 열리는지 확인하는 묶음입니다.\n")
        for name in ["문서-샘플.pdf", "표계산-샘플.xlsx", "이미지-풍경.png", "마크다운-샘플.md"]:
            f = os.path.join(OUT, name)
            if os.path.exists(f):
                z.write(f, "자료/" + name)
        z.writestr("코드/hello.py", 'print("압축 안 파이썬 파일도 실행됩니다")\nfor i in range(5):\n    print(i, i * i)\n')
        z.writestr("코드/도우미.py", 'def 인사(이름):\n    return f"안녕, {이름}!"\n')
        z.writestr("웹/index.html", '<!doctype html><meta charset="utf-8"><h1>압축 안 HTML</h1>'
                                    '<img src="../자료/이미지-풍경.png" width="300">')
    done("압축-샘플.zip")

    with tarfile.open(p("압축-샘플.tgz"), "w:gz") as t:
        for name in ["텍스트-샘플.txt", "코드-샘플.js", "설정-샘플.json", "마크다운-샘플.md"]:
            f = os.path.join(OUT, name)
            if os.path.exists(f):
                t.add(f, arcname="묶음/" + name)
    done("압축-샘플.tgz")


# ---------------------------------------------------------------- 안내 문서
def make_readme():
    write_text("00-읽어보세요.md", '''# 만능파일교실 테스트 파일 모음

이 폴더 하나를 프로그램에 **드래그**하거나 **열기 → 폴더 열기(Ctrl+Shift+O)** 로 열면,
지원하는 형식을 한 번에 확인할 수 있습니다.

> 이 폴더는 `tools/make-sample-files.py` 로 다시 만들 수 있습니다(모두 가짜 데이터라 마음껏 고치고 지워도 됩니다).

## 들어 있는 파일

| 파일 | 형식 | 확인할 것 |
|------|------|-----------|
| `문서-샘플.pdf` | PDF | 서명·텍스트·날짜·체크, 필기, 페이지 정리, 찾기(`검색테스트`) |
| `문서-샘플.docx` | Word | 미리보기(제목·표·긴 문단) |
| `문서-구형샘플.doc` | Word 97 | 글자만 뽑는 간이 미리보기 |
| `표계산-샘플.xlsx` | Excel | 셀 편집·수식·필터·정렬·차트·저장 (3개 시트) |
| `표계산-샘플.csv` | CSV | 대용량 페이지 보기, XLSX 변환 |
| `발표-샘플.pptx` | PowerPoint | 슬라이드 미리보기, EXE+PPT 설치 시 PDF 정확 변환 |
| `한글-샘플.hwpx` | 한글(신형) | 문단·표·서식 간이 미리보기 |
| `이미지-*.png/.jpg/.webp/.bmp/.gif/.ico/.svg/.avif` | 이미지 | 회전·자르기·모자이크·보정·저장, 확대/축소 |
| `사진모음/` | 폴더 | 폴더 안 이미지 격자 모아보기 |
| `영상-샘플.mp4` | 영상 | 재생, 장면 캡처, 자막 자동 연결(`영상-샘플.srt`) |
| `영상-웹형식.webm` | 영상 | 재생, 자막 자동 연결(`영상-웹형식.vtt`) |
| `영상-변환필요.avi` | 영상 | 브라우저가 못 여는 형식 → EXE 의 MP4 변환 안내 확인 |
| `자막-구형.smi` | 자막 | 수동 연결·표시/숨기기·글자 크기 |
| `오디오-샘플.wav` | 오디오 | 재생(도·미·솔·도 소리가 납니다) |
| `오디오-무음.mp3` | 오디오 | 재생 컨트롤 확인 (소리는 없는 무음 파일) |
| `파이썬-계산기.py` | 파이썬 | 실행·표준입력(`90 85 77` 입력) |
| `파이썬-그래프.py` | 파이썬 | matplotlib 그래프 출력 |
| `파이썬-오류.py` | 파이썬 | 오류(traceback) 표시·실시간 진단 |
| `파이썬프로젝트/` | 폴더 | 옆 파일 자동완성·자동 import·정의 이동(Ctrl+클릭) |
| `노트북-샘플.ipynb` | 주피터 | 셀 편집·실행·순서 바꾸기·커널 상태 유지 |
| `데이터베이스-샘플.db` | SQLite | 테이블·뷰·행 읽기 (EXE + 로컬 파이썬) |
| `마크다운-샘플.md` / `웹페이지-샘플.html` / `텍스트-샘플.txt` | 문서 | 렌더·격리 미리보기·텍스트 보기 |
| `코드-샘플.js/.ts`, `설정-샘플.json/.yaml`, `스타일-샘플.css`, `쿼리-샘플.sql`, `데이터-샘플.xml` | 코드 | 구문 강조·줄 번호 |
| `압축-샘플.zip` / `압축-샘플.tgz` | 압축 | 폴더처럼 열기, 압축 안 파이썬 실행 |
| `수업리플레이-샘플.lesson` | 리플레이 | 판서 되감아 재생·속도 조절 |
| `과제-샘플.task` | 과제 | 문제 보기 → 코드 작성 → ✓ 채점 → 제출본 내보내기 |
| `블록문서-샘플.mnote` | 블록 문서 | 글·표·이미지 블록 편집, HTML/Markdown 내보내기 |

## 통합 검색 확인

여러 파일에 검색용 낱말을 심어 두었습니다.
`검색테스트` 로 검색하면 PDF·워드·엑셀·마크다운·텍스트·HTML 등 여러 파일이 함께 걸립니다.

## 이 폴더에 없는 형식

인코더나 원본 없이 만들 수 없는 형식이라 빠졌습니다. 필요하면 실제 파일을 이 폴더에 넣어 함께 확인하세요.

- `.hwp`(구형 한글), `.xls`(구형 엑셀) — 한글/엑셀에서 직접 저장해야 하는 바이너리 형식
- `.m4a` `.aac` `.ogg` `.flac` — 오디오 인코더가 필요합니다
- `.mov` `.mkv` `.wmv` `.flv` — 영상 인코더가 필요합니다 (대신 `.avi` 로 변환 안내를 확인할 수 있습니다)

## 참고

- `과제-샘플.task` 의 정답은 `합계 = sum(수들)`, `평균 = 합계 / len(수들)` 입니다.
- `오디오-무음.mp3` 는 인코더 없이 만든 무음 파일이라 소리가 나지 않습니다(재생/탐색 UI 확인용).
''')
    done("00-읽어보세요.md")


# ---------------------------------------------------------------- 실행
def main():
    os.makedirs(OUT, exist_ok=True)
    print("생성 위치:", OUT)
    step("PDF", make_pdf)
    step("DOCX", make_docx)
    step("구형 DOC(node 필요)", make_doc_legacy)
    step("XLSX·CSV", make_xlsx)
    step("PPTX", make_pptx)
    step("HWPX", make_hwpx)
    step("이미지", make_images)
    step("영상", make_video)
    step("자막", make_subtitles)
    step("오디오", make_audio)
    step("파이썬", make_python)
    step("주피터 노트북", make_notebook)
    step("SQLite", make_db)
    step("텍스트·코드·설정", make_texts)
    step("수업 리플레이(.lesson)", make_lesson)
    step("과제 패키지(.task)", make_task)
    step("블록 문서(.mnote)", make_mnote)
    step("압축", make_archives)
    step("안내 문서", make_readme)
    print("\n완료: %d 항목 생성" % len(made))
    if skipped:
        print("건너뜀 %d 개:" % len(skipped))
        for s in skipped:
            print("  -", s)


if __name__ == "__main__":
    main()
